"""
OCR Service - Handles text extraction from documents.
Supports: EasyOCR (free), Azure Document Intelligence (premium), GPT-4V
"""

import asyncio
import os
import tempfile
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
import structlog

from app.config import OCRProvider, get_settings

logger = structlog.get_logger()


class BaseOCR(ABC):
    """Abstract base class for OCR providers."""
    
    @abstractmethod
    async def extract_text(self, file_path: str) -> tuple[str, int]:
        """Extract text from document. Returns (text, page_count)."""
        pass
    
    @abstractmethod
    def is_available(self) -> bool:
        """Check if this OCR provider is available."""
        pass


class EasyOCRProvider(BaseOCR):
    """
    EasyOCR - Free, runs locally on CPU/GPU.
    More Windows-friendly than PaddleOCR.
    Accuracy: 90%+ for clear documents.
    """
    
    def __init__(self):
        self._reader = None
        self._initialized = False
    
    def _initialize(self, languages: list[str] = None):
        """Lazy initialization of EasyOCR."""
        if not self._initialized:
            try:
                import easyocr
                # Default to English, can be extended
                langs = languages or ['en']
                logger.info("Loading EasyOCR model...", languages=langs)
                self._reader = easyocr.Reader(
                    langs,
                    gpu=False,  # Use CPU by default for compatibility
                    verbose=False
                )
                self._initialized = True
                logger.info("EasyOCR loaded successfully")
            except Exception as e:
                logger.error("Failed to initialize EasyOCR", error=str(e))
                self._initialized = False
    
    def is_available(self) -> bool:
        try:
            import easyocr
            return True
        except ImportError:
            return False
    
    async def extract_text(self, file_path: str, languages: list[str] = None) -> tuple[str, int]:
        """Extract text using EasyOCR."""
        self._initialize(languages)
        if not self._reader:
            raise RuntimeError("EasyOCR not initialized")
        
        file_path = Path(file_path)
        all_text = []
        page_count = 0
        
        # Handle different file types
        if file_path.suffix.lower() == '.pdf':
            # Convert PDF pages to images
            doc = fitz.open(str(file_path))
            page_count = len(doc)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                # Render at 200 DPI for good OCR quality
                pix = page.get_pixmap(matrix=fitz.Matrix(200/72, 200/72))
                img_data = pix.tobytes("png")
                
                # Save temp image and run OCR
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    tmp.write(img_data)
                    tmp_path = tmp.name
                
                try:
                    # Run OCR in thread pool to not block
                    result = await asyncio.get_event_loop().run_in_executor(
                        None, self._reader.readtext, tmp_path
                    )
                    
                    # Extract text from results
                    page_text = [detection[1] for detection in result]
                    all_text.append(' '.join(page_text))
                    
                    logger.info(f"Processed page {page_num + 1}/{page_count}")
                finally:
                    os.unlink(tmp_path)
            
            doc.close()
        
        elif file_path.suffix.lower() in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
            page_count = 1
            result = await asyncio.get_event_loop().run_in_executor(
                None, self._reader.readtext, str(file_path)
            )
            page_text = [detection[1] for detection in result]
            all_text.append(' '.join(page_text))
        
        else:
            # For other documents, try to extract text directly first
            text = await self._extract_from_office(file_path)
            if text:
                return text, 1
        
        return '\n\n'.join(all_text), page_count
    
    async def _extract_from_office(self, file_path: Path) -> Optional[str]:
        """Extract text from Office documents."""
        suffix = file_path.suffix.lower()
        
        try:
            if suffix == '.docx':
                from docx import Document
                doc = Document(str(file_path))
                return '\n'.join([para.text for para in doc.paragraphs])
            
            elif suffix == '.xlsx':
                from openpyxl import load_workbook
                wb = load_workbook(str(file_path))
                text_parts = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = [str(cell) for cell in row if cell]
                        if row_text:
                            text_parts.append('\t'.join(row_text))
                return '\n'.join(text_parts)
            
            elif suffix == '.pptx':
                from pptx import Presentation
                prs = Presentation(str(file_path))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text_parts.append(shape.text)
                return '\n'.join(text_parts)
        
        except Exception as e:
            logger.warning("Office extraction failed", error=str(e))
        
        return None


class AzureOCRProvider(BaseOCR):
    """
    Azure Document Intelligence - Premium, highest accuracy (97%+).
    Cost: ~$0.01 per page.
    """
    
    def __init__(self):
        self.settings = get_settings()
        self._client = None
    
    def is_available(self) -> bool:
        return bool(self.settings.azure_doc_endpoint and self.settings.azure_doc_key)
    
    def _get_client(self):
        if not self._client:
            try:
                from azure.ai.documentintelligence import DocumentIntelligenceClient
                from azure.core.credentials import AzureKeyCredential
                
                self._client = DocumentIntelligenceClient(
                    endpoint=self.settings.azure_doc_endpoint,
                    credential=AzureKeyCredential(self.settings.azure_doc_key)
                )
            except ImportError:
                logger.warning("Azure SDK not installed. Run: pip install azure-ai-documentintelligence")
                return None
        return self._client
    
    async def extract_text(self, file_path: str) -> tuple[str, int]:
        """Extract text using Azure Document Intelligence."""
        client = self._get_client()
        if not client:
            raise RuntimeError("Azure client not available")
        
        with open(file_path, 'rb') as f:
            document_bytes = f.read()
        
        # Run in thread pool
        def analyze():
            poller = client.begin_analyze_document(
                "prebuilt-read",
                document_bytes,
                content_type="application/octet-stream"
            )
            return poller.result()
        
        result = await asyncio.get_event_loop().run_in_executor(None, analyze)
        
        # Extract text from result
        all_text = []
        page_count = len(result.pages) if result.pages else 0
        
        for page in result.pages:
            page_text = []
            for line in page.lines:
                page_text.append(line.content)
            all_text.append('\n'.join(page_text))
        
        return '\n\n'.join(all_text), page_count


class GPT4OCRProvider(BaseOCR):
    """
    GPT-4 Vision - Alternative premium option.
    Good for mixed content (text + diagrams).
    Cost: ~$0.02 per page.
    """
    
    def __init__(self):
        self.settings = get_settings()
    
    def is_available(self) -> bool:
        return bool(self.settings.openai_api_key)
    
    async def extract_text(self, file_path: str) -> tuple[str, int]:
        """Extract text using GPT-4 Vision."""
        import base64
        from openai import AsyncOpenAI
        
        client = AsyncOpenAI(api_key=self.settings.openai_api_key)
        file_path = Path(file_path)
        all_text = []
        page_count = 0
        
        if file_path.suffix.lower() == '.pdf':
            doc = fitz.open(str(file_path))
            page_count = len(doc)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(matrix=fitz.Matrix(150/72, 150/72))
                img_bytes = pix.tobytes("png")
                img_base64 = base64.b64encode(img_bytes).decode()
                
                response = await client.chat.completions.create(
                    model="gpt-4o",
                    messages=[{
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Extract ALL text from this document image. Preserve the original structure and formatting. Return only the extracted text, nothing else."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{img_base64}"
                                }
                            }
                        ]
                    }],
                    max_tokens=4096
                )
                all_text.append(response.choices[0].message.content)
                logger.info(f"Processed page {page_num + 1}/{page_count}")
            
            doc.close()
        else:
            # Handle single images
            page_count = 1
            with open(file_path, "rb") as f:
                img_bytes = f.read()
            img_base64 = base64.b64encode(img_bytes).decode()
            
            response = await client.chat.completions.create(
                model="gpt-4o",
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Extract ALL text from this image. Return only the extracted text."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{img_base64}"
                            }
                        }
                    ]
                }],
                max_tokens=4096
            )
            all_text.append(response.choices[0].message.content)
        
        return '\n\n'.join(all_text), page_count


class OCRService:
    """Main OCR service that manages different providers."""
    
    def __init__(self):
        self.providers = {
            OCRProvider.PADDLEOCR: EasyOCRProvider(),  # Using EasyOCR as PADDLEOCR replacement
            OCRProvider.AZURE: AzureOCRProvider(),
            OCRProvider.GPT4O: GPT4OCRProvider(),
        }
    
    def get_available_providers(self) -> list[str]:
        """Get list of available OCR providers."""
        available = []
        for p, impl in self.providers.items():
            if impl.is_available():
                # Map back to user-friendly name
                name = "easyocr" if p == OCRProvider.PADDLEOCR else p.value
                available.append(name)
        return available
    
    async def extract_text(
        self, 
        file_path: str, 
        provider: Optional[OCRProvider] = None
    ) -> tuple[str, int, str]:
        """
        Extract text from document.
        Returns: (extracted_text, page_count, provider_used)
        """
        # Default to EasyOCR if no provider specified
        if provider is None:
            settings = get_settings()
            if settings.translation_mode.value == "premium" and self.providers[OCRProvider.AZURE].is_available():
                provider = OCRProvider.AZURE
            else:
                provider = OCRProvider.PADDLEOCR  # Uses EasyOCR
        
        ocr = self.providers.get(provider)
        if not ocr or not ocr.is_available():
            # Fallback to EasyOCR
            logger.warning(f"Provider {provider} not available, falling back to EasyOCR")
            ocr = self.providers[OCRProvider.PADDLEOCR]
            provider = OCRProvider.PADDLEOCR
        
        provider_name = "easyocr" if provider == OCRProvider.PADDLEOCR else provider.value
        logger.info("Starting OCR extraction", provider=provider_name, file=file_path)
        text, page_count = await ocr.extract_text(file_path)
        logger.info("OCR extraction complete", pages=page_count, chars=len(text))
        
        return text, page_count, provider_name
